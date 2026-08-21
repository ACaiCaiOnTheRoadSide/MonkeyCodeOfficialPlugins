#!/usr/bin/env python3
"""Render each HTML section.slide with Chromium and package PNGs as a 16:9 PPTX."""

from __future__ import annotations

import argparse
import hashlib
import importlib.util
import json
import os
import sys
import tempfile
from pathlib import Path

SLIDE_WIDTH_PX = 1600
SLIDE_HEIGHT_PX = 900


def fail(message: str) -> "None":
    print(f"error: {message}", file=sys.stderr)
    raise SystemExit(2)


def require_dependencies() -> None:
    missing = []
    if importlib.util.find_spec("playwright") is None:
        missing.append("playwright (python -m pip install playwright)")
    if importlib.util.find_spec("pptx") is None:
        missing.append("python-pptx (python -m pip install python-pptx)")
    if missing:
        fail("missing Python dependencies: " + ", ".join(missing) + ". This script does not install dependencies.")


def commit_outputs(presentation: object, manifest: dict[str, object], output: Path) -> None:
    manifest_output = output.with_suffix(output.suffix + ".fidelity.json")
    for target in (output, manifest_output):
        if target.is_dir():
            fail(f"output target is a directory: {target}")

    with tempfile.TemporaryDirectory(prefix=".monkeycode-slides-", dir=output.parent) as staging_dir:
        staging = Path(staging_dir)
        staged_pptx = staging / "deck.pptx"
        staged_manifest = staging / "deck.pptx.fidelity.json"
        presentation.save(staged_pptx)
        committed_manifest = dict(manifest)
        committed_manifest["pptx_sha256"] = hashlib.sha256(staged_pptx.read_bytes()).hexdigest()
        staged_manifest.write_text(json.dumps(committed_manifest, indent=2) + "\n", encoding="utf-8")

        backups: dict[Path, Path] = {}
        installed: list[Path] = []
        try:
            for index, target in enumerate((output, manifest_output)):
                if target.exists() or target.is_symlink():
                    backup = staging / f"backup-{index}"
                    os.replace(target, backup)
                    backups[target] = backup
            for staged, target in ((staged_pptx, output), (staged_manifest, manifest_output)):
                os.replace(staged, target)
                installed.append(target)
        except OSError as exc:
            rollback_errors = []
            for target in reversed(installed):
                try:
                    target.unlink(missing_ok=True)
                except OSError as rollback_exc:
                    rollback_errors.append(str(rollback_exc))
            for target, backup in backups.items():
                try:
                    os.replace(backup, target)
                except OSError as rollback_exc:
                    rollback_errors.append(str(rollback_exc))
            detail = f"; rollback errors: {', '.join(rollback_errors)}" if rollback_errors else ""
            fail(f"could not atomically install PPTX and fidelity manifest: {exc}{detail}")


def export(source: Path, output: Path) -> int:
    require_dependencies()
    from playwright.sync_api import Error as PlaywrightError
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Inches

    if not source.is_file():
        fail(f"HTML input does not exist: {source}")
    output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="monkeycode-slides-") as temp_dir:
        screenshots: list[Path] = []
        try:
            with sync_playwright() as playwright:
                try:
                    browser = playwright.chromium.launch(headless=True)
                except PlaywrightError as exc:
                    fail(
                        "Chromium is unavailable. Install the Playwright browser explicitly with "
                        f"`python -m playwright install chromium`. Original error: {exc}"
                    )
                page = browser.new_page(
                    viewport={"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
                    device_scale_factor=1,
                )
                page.goto(source.resolve().as_uri(), wait_until="networkidle")
                page.add_style_tag(content="""
                    html, body { margin: 0 !important; padding: 0 !important; background: transparent !important; }
                    *, *::before, *::after {
                      animation: none !important; transition: none !important;
                      caret-color: transparent !important;
                    }
                    section.slide {
                      box-sizing: border-box !important;
                      width: 1600px !important; height: 900px !important;
                      min-width: 1600px !important; min-height: 900px !important;
                      max-width: 1600px !important; max-height: 900px !important;
                      margin: 0 !important; transform: none !important;
                    }
                """)
                page.evaluate("document.fonts && document.fonts.ready")
                slides = page.locator("section.slide")
                count = slides.count()
                if count == 0:
                    browser.close()
                    fail('no slides found; expected one or more `<section class="slide">` elements')
                for index in range(count):
                    slide = slides.nth(index)
                    slide.scroll_into_view_if_needed()
                    box = slide.bounding_box()
                    if not box:
                        fail(f"slide {index + 1} is not visible and cannot be rendered")
                    if round(box["width"]) != SLIDE_WIDTH_PX or round(box["height"]) != SLIDE_HEIGHT_PX:
                        fail(f"slide {index + 1} did not resolve to 1600x900 (got {box['width']}x{box['height']})")
                    screenshot = Path(temp_dir) / f"slide-{index + 1:03d}.png"
                    slide.screenshot(path=str(screenshot), animations="disabled")
                    screenshots.append(screenshot)
                browser.close()
        except PlaywrightError as exc:
            fail(f"Chromium could not load or render {source}: {exc}")

        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        blank = presentation.slide_layouts[6]
        slide_hashes = []
        for index, screenshot in enumerate(screenshots, 1):
            slide = presentation.slides.add_slide(blank)
            picture = slide.shapes.add_picture(
                str(screenshot), 0, 0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
            picture.name = "full-slide-render"
            slide_hashes.append({"index": index, "sha256": hashlib.sha256(screenshot.read_bytes()).hexdigest()})
        manifest = {
            "mode": "high-fidelity-image",
            "editable": False,
            "source": str(source.resolve()),
            "slide_count": len(screenshots),
            "canvas_px": [SLIDE_WIDTH_PX, SLIDE_HEIGHT_PX],
            "slides": slide_hashes,
        }
        commit_outputs(presentation, manifest, output)
    return len(screenshots)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("html", type=Path, help="source HTML deck")
    parser.add_argument("pptx", type=Path, help="output PPTX path")
    args = parser.parse_args()
    if args.pptx.suffix.lower() != ".pptx":
        fail("output path must end in .pptx")
    count = export(args.html, args.pptx)
    print(
        f"exported {count} slide(s) to {args.pptx} in high-fidelity image mode (non-editable); "
        f"wrote {args.pptx.with_suffix(args.pptx.suffix + '.fidelity.json')}"
    )


if __name__ == "__main__":
    main()
